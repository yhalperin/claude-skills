"""
L3-Agents Status PPTX Generator  (v2 — health indicators)
Usage: python generate_pptx.py --data <path-to-data.json>

Input JSON schema:
{
  "rows": [
    {
      "agent":               "Agent Name (KEY)" | "Not mapped yet",
      "theme":               "Theme Summary\\n(IAI-XXXX)",
      "business_value":      "one-line value",
      "impact":              "one-line impact",
      "status":              "In Progress" | "Planned" | "Open" | ...,
      "time_in_status_days": 14,
      "pi":                  "27-Q1",
      "finish_date":         "2026-10-15" | null,
      "health":              "On Track" | "At Risk" | "Off Track"
    }
  ],
  "output_dir": "C:\\path\\to\\output"
}

Columns: Agent | Theme | Business Value | Impact | Status (+time) | ETA | Health | Dev Phase removed
Output:  <output_dir>/l3-agents-status-YYYY-MM-DD.pptx
"""

import argparse
import json
import os
import sys
from datetime import date, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Column headers & widths ────────────────────────────────────────────────────
# Total must ≈ slide_width − 2 × margin = 13.33 − 0.76 = 12.57 inches

HEADERS = [
    "Agent (Master Feature)",
    "Theme Name (Jira ID)",
    "Business Value",
    "Impact",
    "Status",
    "ETA",
    "Health",
]

# inches; sum = 12.58
COL_W = [1.90, 2.00, 2.35, 1.90, 1.25, 0.83, 1.35]

# ── Color palette ─────────────────────────────────────────────────────────────

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG     = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG    = RGBColor(0x1F, 0x6B, 0xC8)
HEADER_TXT   = WHITE
ROW_BG       = WHITE
ROW_ALT_BG   = RGBColor(0xEB, 0xF3, 0xFB)
BORDER_COLOR = RGBColor(0xC0, 0xCC, 0xD8)
TITLE_COLOR  = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_CLR = RGBColor(0x55, 0x55, 0x66)
DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x1A)
DIM_TEXT     = RGBColor(0x55, 0x55, 0x66)

# Jira-aligned status colors
STATUS_COLORS = {
    "Done":                 RGBColor(0x42, 0x52, 0x6E),
    "In Progress":          RGBColor(0x00, 0x52, 0xCC),
    "HL Dev Discovery":     RGBColor(0x42, 0x52, 0x6E),
    "HL Product Discovery": RGBColor(0x42, 0x52, 0x6E),
    "Planned":              RGBColor(0x21, 0x6E, 0x4E),
    "Open":                 RGBColor(0x42, 0x52, 0x6E),
}

# Health — background fill + text color pairs
HEALTH_BG = {
    "On Track":  RGBColor(0xD6, 0xF0, 0xDE),   # light green
    "At Risk":   RGBColor(0xFF, 0xEE, 0xC2),    # light amber
    "Off Track": RGBColor(0xFC, 0xD9, 0xD9),    # light red
}
HEALTH_TXT = {
    "On Track":  RGBColor(0x1A, 0x60, 0x35),   # dark green
    "At Risk":   RGBColor(0x7A, 0x50, 0x00),    # dark amber
    "Off Track": RGBColor(0x99, 0x1B, 0x1B),    # dark red
}
HEALTH_LABELS = {
    "On Track":  "On Track",
    "At Risk":   "At Risk",
    "Off Track": "Off Track",
}

# Status legend for top-right of slide
STATUS_LEGEND = [
    (RGBColor(0x00, 0x52, 0xCC), "In Progress"),
    (RGBColor(0x21, 0x6E, 0x4E), "Planned"),
    (RGBColor(0x42, 0x52, 0x6E), "Open / Discovery / Other"),
]

# Health legend
HEALTH_LEGEND = [
    (RGBColor(0x1A, 0x60, 0x35), "On Track"),
    (RGBColor(0x7A, 0x50, 0x00), "At Risk"),
    (RGBColor(0x99, 0x1B, 0x1B), "Off Track"),
]

# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_cell_bg(cell, rgb: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _add_cell_border(cell, rgb: RGBColor, width_pt: float = 0.5):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w    = int(width_pt * 12700)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set("w", str(w))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _set_cell_text(cell, lines, font_sizes, bolds, colors, align=PP_ALIGN.LEFT):
    """
    lines      – list of str (one per paragraph)
    font_sizes – list of float (pt) or single float
    bolds      – list of bool or single bool
    colors     – list of RGBColor or single RGBColor
    """
    if isinstance(lines, str):
        lines = [lines]
    if not isinstance(font_sizes, list):
        font_sizes = [font_sizes] * len(lines)
    if not isinstance(bolds, list):
        bolds = [bolds] * len(lines)
    if not isinstance(colors, list):
        colors = [colors] * len(lines)

    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (line, fs, bold, color) in enumerate(zip(lines, font_sizes, bolds, colors)):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text            = line
        run.font.size       = Pt(fs)
        run.font.bold       = bold
        run.font.color.rgb  = color


def _add_dot(slide, left, top, size, rgb: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, size, size)
    from pptx.oxml.ns import qn
    sp    = shape._element
    spPr  = sp.find(qn("p:spPr"))
    geom  = spPr.find(qn("a:prstGeom"))
    if geom is not None:
        geom.set("prst", "ellipse")
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _format_eta(finish_date: str | None) -> str:
    if not finish_date:
        return "—"
    try:
        d = datetime.strptime(finish_date, "%Y-%m-%d").date()
        return d.strftime("%b %-d")   # e.g. "Oct 15"
    except Exception:
        try:
            d = datetime.strptime(finish_date, "%Y-%m-%d").date()
            return d.strftime("%b %d")
        except Exception:
            return finish_date[:10]


def _format_time_in_status(days: int | None) -> str:
    if days is None:
        return ""
    if days == 0:
        return "today"
    if days < 7:
        return f"{days}d"
    weeks = days // 7
    rem   = days % 7
    return f"{weeks}w {rem}d" if rem else f"{weeks}w"


# ── Slide builder ─────────────────────────────────────────────────────────────

def add_slide(prs, chunk, slide_num, total_slides, today_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLIDE_BG

    SW = prs.slide_width
    SH = prs.slide_height
    M  = Inches(0.38)

    # ── Title ─────────────────────────────────────────────────────────────────
    tb  = slide.shapes.add_textbox(M, M, Inches(7.5), Inches(0.50))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text           = "L3 Agents Progress"
    run.font.size      = Pt(22)
    run.font.bold      = True
    run.font.color.rgb = TITLE_COLOR

    sb  = slide.shapes.add_textbox(M, Inches(0.86), Inches(9.0), Inches(0.22))
    run2 = sb.text_frame.paragraphs[0].add_run()
    run2.text = (
        f"Project IAI · label: L3-Agents · ETA: cf[21221] · Agent: cf[11140] · "
        f"{today_str}  —  Slide {slide_num} of {total_slides}"
    )
    run2.font.size      = Pt(8)
    run2.font.color.rgb = SUBTITLE_CLR

    # ── Status legend (top-right) ─────────────────────────────────────────────
    dot_r   = Inches(0.12)
    dot_gap = Inches(0.05)
    lbl_w   = Inches(1.40)
    leg_top = Inches(0.22)
    n_items = len(STATUS_LEGEND)
    total_w = n_items * (dot_r + dot_gap + lbl_w + Inches(0.08))
    leg_x   = SW - M - total_w + Inches(0.08)

    for i, (color, label) in enumerate(STATUS_LEGEND):
        x = leg_x + i * (dot_r + dot_gap + lbl_w + Inches(0.08))
        _add_dot(slide, x, leg_top + Inches(0.04), dot_r, color)
        ltb = slide.shapes.add_textbox(x + dot_r + dot_gap, leg_top, lbl_w, Inches(0.22))
        lr  = ltb.text_frame.paragraphs[0].add_run()
        lr.text           = label
        lr.font.size      = Pt(8)
        lr.font.color.rgb = DARK_TEXT

    # ── Health legend (second row, top-right) ─────────────────────────────────
    hlth_lbl_w = Inches(0.85)
    hlth_top   = Inches(0.46)
    n_hlth     = len(HEALTH_LEGEND)
    total_hw   = n_hlth * (dot_r + dot_gap + hlth_lbl_w + Inches(0.08))
    hlth_x     = SW - M - total_hw + Inches(0.08)

    for i, (color, label) in enumerate(HEALTH_LEGEND):
        x = hlth_x + i * (dot_r + dot_gap + hlth_lbl_w + Inches(0.08))
        _add_dot(slide, x, hlth_top + Inches(0.04), dot_r, color)
        ltb = slide.shapes.add_textbox(x + dot_r + dot_gap, hlth_top, hlth_lbl_w, Inches(0.22))
        lr  = ltb.text_frame.paragraphs[0].add_run()
        lr.text           = label
        lr.font.size      = Pt(8)
        lr.font.color.rgb = DARK_TEXT

    # ── Table ─────────────────────────────────────────────────────────────────
    n_rows     = len(chunk) + 1
    table_top  = Inches(1.10)
    table_h    = SH - table_top - M - Inches(0.18)
    table_w    = SW - 2 * M
    row_h_hdr  = Inches(0.38)
    row_h_data = (table_h - row_h_hdr) / max(len(chunk), 1)

    tbl_shape = slide.shapes.add_table(n_rows, len(HEADERS), M, table_top, table_w, table_h)
    tbl = tbl_shape.table

    total_col_w = sum(COL_W)
    for ci, w in enumerate(COL_W):
        tbl.columns[ci].width = int(table_w * w / total_col_w)

    tbl.rows[0].height = row_h_hdr
    for ri in range(1, n_rows):
        tbl.rows[ri].height = int(row_h_data)

    # Header row
    for ci, hdr in enumerate(HEADERS):
        cell = tbl.cell(0, ci)
        _set_cell_bg(cell, HEADER_BG)
        _add_cell_border(cell, BORDER_COLOR, 0.5)
        _set_cell_text(cell, hdr, 10, True, HEADER_TXT)
        cell.margin_top = cell.margin_bottom = Pt(5)
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(4)

    # Data rows
    for ri, row in enumerate(chunk):
        bg_row   = ROW_ALT_BG if ri % 2 == 1 else ROW_BG
        s_color  = STATUS_COLORS.get(row["status"], DIM_TEXT)
        health   = row.get("health", "On Track")
        h_bg     = HEALTH_BG.get(health, ROW_BG)
        h_txt    = HEALTH_TXT.get(health, DARK_TEXT)

        tis_days = row.get("time_in_status_days")
        tis_str  = _format_time_in_status(tis_days)
        eta_str  = _format_eta(row.get("finish_date"))

        # column index → (lines, sizes, bolds, colors, bg_override)
        col_data = [
            # Agent
            ([row["agent"]], [8.5], [False], [DIM_TEXT], bg_row),
            # Theme
            ([row["theme"]], [8.5], [True],  [DARK_TEXT], bg_row),
            # Business Value
            ([row["business_value"]], [8.5], [False], [DARK_TEXT], bg_row),
            # Impact
            ([row["impact"]], [8.5], [False], [DIM_TEXT], bg_row),
            # Status (+ time in status)
            (
                [row["status"], tis_str] if tis_str else [row["status"]],
                [9, 7.5],
                [True, False],
                [s_color, DIM_TEXT],
                bg_row,
            ),
            # ETA
            ([eta_str], [8.5], [False], [DARK_TEXT], bg_row),
            # Health — gets its own background color
            ([HEALTH_LABELS[health]], [8.5], [True], [h_txt], h_bg),
        ]

        for ci, (lines, sizes, bolds, colors, bg_ci) in enumerate(col_data):
            cell = tbl.cell(ri + 1, ci)
            _set_cell_bg(cell, bg_ci)
            _add_cell_border(cell, BORDER_COLOR, 0.4)
            _set_cell_text(cell, lines, sizes, bolds, colors)
            cell.margin_top = cell.margin_bottom = Pt(4)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(4)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Generate L3-Agents status PPTX")
    parser.add_argument("--data", required=True, help="Path to JSON data file")
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    rows       = data["rows"]
    output_dir = data.get("output_dir", r"C:\Users\yhalperin\Documents\L3_status_presentations")

    os.makedirs(output_dir, exist_ok=True)

    today_str = date.today().strftime("%Y-%m-%d")
    out_file  = os.path.join(output_dir, f"l3-agents-status-{today_str}.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    CHUNK_SIZE = 6
    chunks = [rows[i:i + CHUNK_SIZE] for i in range(0, len(rows), CHUNK_SIZE)]
    if not chunks:
        print("ERROR: No rows provided", file=sys.stderr)
        sys.exit(1)

    for idx, chunk in enumerate(chunks):
        add_slide(prs, chunk, idx + 1, len(chunks), today_str)

    prs.save(out_file)
    print(out_file)

    on_track  = sum(1 for r in rows if r.get("health") == "On Track")
    at_risk   = sum(1 for r in rows if r.get("health") == "At Risk")
    off_track = sum(1 for r in rows if r.get("health") == "Off Track")
    print(
        f"Slides: {len(chunks)}  Rows: {len(rows)}  "
        f"Health: {on_track} On Track / {at_risk} At Risk / {off_track} Off Track",
        file=sys.stderr,
    )


if __name__ == "__main__":
    main()
